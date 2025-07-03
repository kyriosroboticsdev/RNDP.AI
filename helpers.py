# helpers.py

from typing import Type, Optional, List
from pydantic import BaseModel
from openai import AzureOpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image
import os
import datetime
from dotenv import load_dotenv
import fitz  # PyMuPDF for PDF parsing

# Load environment variables
load_dotenv()

AZURE_KEY = os.getenv("AZURE_OPENAI_KEY")
AZURE_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT")

if not AZURE_KEY or not AZURE_ENDPOINT or not AZURE_DEPLOYMENT:
    raise ValueError("Azure OpenAI environment variables are not set")

client = AzureOpenAI(
    api_key=AZURE_KEY,
    api_version="2024-02-15-preview",
    azure_endpoint=AZURE_ENDPOINT
)

def get_ai_recommended_image(topic: str, filenames: List[str]) -> str:
    """Select the most relevant image filename based on a topic using Azure OpenAI."""
    list_text = ", ".join(filenames)
    response = client.chat.completions.create(
        model=AZURE_DEPLOYMENT,
        messages=[
            {"role": "system", "content": "You choose relevant image filenames for robotics topics."},
            {"role": "user", "content": f"Which of these best fits the topic '{topic}': {list_text}"}
        ],
        temperature=0.3,
        max_tokens=100
    )
    answer = response.choices[0].message.content.strip()
    for name in filenames:
        if name in answer:
            return name
    return filenames[0] if filenames else ""

def extract_text_from_slide_file(uploaded_file) -> str:
    """Extract text from a .pptx or .pdf file."""
    extension = uploaded_file.name.lower().split(".")[-1]
    text = ""

    if extension == "pptx":
        presentation = Presentation(uploaded_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif extension == "pdf":
        pdf = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page in pdf:
            text += page.get_text()

    return text.strip()

def structured_generator(
    model_name: str,
    prompt: str,
    output_model: Type[BaseModel],
    template_file,
    content_text: str,
    image_file: Optional[BytesIO] = None,
    font_name: str = "Calibri",
    font_color: str = "000000",
    fallback_image_filename: Optional[str] = None
) -> BaseModel:
    """Generate a structured slide based on template, content, and optional image."""
    template_path = "temp_template.pptx"
    with open(template_path, "wb") as f:
        f.write(template_file.read())

    prs = Presentation(template_path)
    slide = prs.slides[0]

    response = client.chat.completions.create(
        model=AZURE_DEPLOYMENT,
        messages=[
            {"role": "system", "content": "You write robotics notebook slide titles and 3 paragraphs of content."},
            {"role": "user", "content": f"Write a short title and a 3+ sentence explanation for: {content_text}"}
        ],
        temperature=0.7,
        max_tokens=600
    )

    result = response.choices[0].message.content.strip()
    title, paragraph = (result.split("\n", 1) if "\n" in result else (content_text, result))

    for shape in slide.shapes:
        if shape.has_text_frame:
            raw_text = shape.text_frame.text.strip().lower()
            if "title" in raw_text:
                shape.text = title
            elif "example text" in raw_text or "content" in raw_text:
                shape.text = paragraph
            elif "date" in raw_text:
                shape.text = datetime.datetime.now().strftime("%B %d, %Y")

            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(18)
                    try:
                        r, g, b = tuple(int(font_color[i:i+2], 16) for i in (0, 2, 4))
                        run.font.color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass

    image_stream = _prepare_image(image_file, fallback_image_filename)
    if image_stream:
        _insert_image(slide, image_stream)

    return _finalize(prs, output_model)

def _prepare_image(image_file: Optional[BytesIO], fallback_filename: Optional[str]) -> Optional[BytesIO]:
    """Prepare image stream from upload or fallback."""
    if image_file:
        return BytesIO(image_file.read())

    if fallback_filename:
        fallback_path = os.path.join("images", fallback_filename)
        if os.path.exists(fallback_path):
            with open(fallback_path, "rb") as f:
                return BytesIO(f.read())

    return None

def _insert_image(slide, image_stream: BytesIO):
    image_stream.seek(0)  # Reset stream position
    img = Image.open(image_stream)
    img_width, img_height = img.size

    for shape in slide.shapes:
        if shape.has_text_frame:
            placeholder_text = " ".join(
                run.text.lower().strip()
                for para in shape.text_frame.paragraphs
                for run in para.runs
            )
            if any(keyword in placeholder_text for keyword in ["insert image here", "image placeholder", "image goes here"]):
                # Save placeholder position and size
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)  # Remove the placeholder textbox

                # Resize image to fit inside the placeholder box
                scale = min(width / img_width, height / img_height)
                new_width = int(img_width * scale)
                new_height = int(img_height * scale)

                image_stream.seek(0)
                slide.shapes.add_picture(
                    image_stream,
                    left + (width - new_width) // 2,
                    top + (height - new_height) // 2,
                    width=new_width,
                    height=new_height
                )
                break

def _finalize(prs: Presentation, output_model: Type[BaseModel]) -> BaseModel:
    """Finalize and serialize the presentation."""
    output_path = "generated_slide.pptx"
    prs.save(output_path)
    with open(output_path, "rb") as f:
        pptx_bytes = f.read()

    class Output(output_model):
        pptx_bytes: bytes

    return Output(pptx_bytes=pptx_bytes)

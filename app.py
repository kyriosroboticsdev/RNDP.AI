# app.py

import streamlit as st
import os
import time
import uuid
import tempfile
from pydantic import BaseModel
from helpers import (
    structured_generator,
    get_ai_recommended_image,
    extract_text_from_slide_file
)
from db import initialize_db, save_slide, get_user_slides, add_user
from google_auth import login_user, logout_user
from pptx import Presentation

# App setup
st.set_page_config(page_title="Robotics Slide Generator", layout="centered")

# Database setup
initialize_db()

# Session timeout config
SESSION_TIMEOUT = 30 * 60
if "login_time" in st.session_state:
    if time.time() - st.session_state["login_time"] > SESSION_TIMEOUT:
        st.session_state.clear()
        st.warning("⏰ Session expired. Please log in again.")
        st.stop()

# Authenticate user
user_info = login_user()
if not user_info:
    st.stop()

# Register user if not in DB
username = user_info["email"].split("@")[0].lower()
add_user(username, user_info["name"], user_info["email"])

# Sidebar user info + logout
st.sidebar.markdown(f"👋 Logged in as: `{user_info['email']}`")
logout_user()

# Main navigation
st.title("🤖 Robotics Slide Generator")
menu = st.selectbox("Choose a page:", ["Generate Slide", "My Slides"])

# Define output data model
class SlideOutput(BaseModel):
    pptx_bytes: bytes

# ===============================
# 🔧 PAGE 1: GENERATE SLIDE
# ===============================
if menu == "Generate Slide":
    st.header("📄 Generate a New Slide")

    input_text = st.text_area("📝 Slide Topic", help="Write a brief prompt or description")
    template_file = st.file_uploader("📂 Upload Template (.pptx)", type=["pptx"])
    image_file = st.file_uploader("🖼️ Optional image for slide", type=["png", "jpg", "jpeg"])

    # Font and image customization
    fallback_images = os.listdir("images") if os.path.exists("images") else []
    fallback_image_choice = None
    if not image_file and fallback_images:
        ai_choice = get_ai_recommended_image(input_text, fallback_images)
        fallback_image_choice = st.selectbox("📎 Choose fallback image", fallback_images, index=fallback_images.index(ai_choice))

    font_name = st.selectbox("🖋 Font Style", ["Calibri", "Arial", "Times New Roman", "Verdana"])
    font_color = st.color_picker("🎨 Font Color", "#000000")[1:]

    if st.button("🚀 Generate Slide"):
        if not input_text or not template_file:
            st.warning("Please provide both a prompt and a template.")
            st.stop()

        result = structured_generator(
            model_name="gpt-4",
            prompt="Generate robotics notebook content.",
            output_model=SlideOutput,
            template_file=template_file,
            content_text=input_text,
            image_file=image_file,
            font_name=font_name,
            font_color=font_color,
            fallback_image_filename=fallback_image_choice
        )

        os.makedirs("slides", exist_ok=True)
        slide_filename = f"slides/{username}_page{uuid.uuid4().hex[:4]}.pptx"
        with open(slide_filename, "wb") as f:
            f.write(result.pptx_bytes)

        title = input_text.split("\n")[0][:40]
        save_slide(username, title, slide_filename)

        st.success("✅ Slide generated and saved!")
        st.download_button("⬇️ Download Slide", result.pptx_bytes, file_name=os.path.basename(slide_filename))

# ===============================
# 📚 PAGE 2: MY SLIDES
# ===============================
elif menu == "My Slides":
    st.header("📚 My Slides")

    slides = get_user_slides(username)
    if not slides:
        st.info("No saved slides yet.")
        st.stop()

    for title, date_created, path in slides:
        st.markdown(f"### 📄 {title}")
        st.caption(f"Created on {date_created}")

        try:
            with open(path, "rb") as f:
                pptx_data = f.read()

            # Save to a temp file and extract text
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(pptx_data)
                tmp.flush()

                prs = Presentation(tmp.name)
                for i, slide in enumerate(prs.slides):
                    st.markdown(f"**Slide {i+1} Content Preview:**")
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text.strip() + "\n"
                    if slide_text:
                        st.text(slide_text.strip())
                    else:
                        st.caption("_No text found on this slide._")
                    break  # Only show first slide preview

            with open(path, "rb") as f:
                st.download_button("⬇️ Download Slide", f, file_name=os.path.basename(path))

        except Exception as e:
            st.warning(f"⚠️ Could not preview slide: {e}")
